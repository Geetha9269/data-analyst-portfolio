from groq import Groq
import pandas as pd
import os
from dotenv import load_dotenv
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from datetime import datetime

# ---- Setup ----
load_dotenv()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# ---- Load Data ----
df = pd.read_csv('data/HR-Em.csv')

print("=" * 50)
print("PHASE 3 - AUTOMATED REPORT GENERATION")
print("=" * 50)

# ---- Calculate All Metrics ----
attrition_rate = round(
    df['Attrition'].value_counts()['Yes'] / len(df) * 100, 2
)

dept_attrition = df.groupby('Department')['Attrition'].apply(
    lambda x: round((x == 'Yes').sum() / len(x) * 100, 2)
)

avg_income_left = round(
    df[df['Attrition'] == 'Yes']['MonthlyIncome'].mean(), 0
)
avg_income_stayed = round(
    df[df['Attrition'] == 'No']['MonthlyIncome'].mean(), 0
)

overtime_attrition = df.groupby('OverTime')['Attrition'].apply(
    lambda x: round((x == 'Yes').sum() / len(x) * 100, 2)
)

top_roles = df.groupby('JobRole')['Attrition'].apply(
    lambda x: round((x == 'Yes').sum() / len(x) * 100, 2)
).sort_values(ascending=False).head(5)

print("✅ Metrics calculated")

# ---- Get AI Insights ----
print("🤖 Getting AI insights...")

response = client.chat.completions.create(
    model="llama-3.3-70b-versatile",
    messages=[
        {
            "role": "system",
            "content": """You are a senior HR Data Analyst 
            writing an executive report. Be professional, 
            specific, and concise."""
        },
        {
            "role": "user",
            "content": f"""
            Write an executive HR Attrition Report with:

            DATA:
            - Total Employees: {len(df)}
            - Attrition Rate: {attrition_rate}%
            - Avg Income Left: ${avg_income_left}
            - Avg Income Stayed: ${avg_income_stayed}
            - Overtime Attrition: {overtime_attrition['Yes']}%
            - No Overtime Attrition: {overtime_attrition['No']}%
            - Sales Dept Attrition: {dept_attrition['Sales']}%
            - R&D Dept Attrition: {dept_attrition['Research & Development']}%
            - HR Dept Attrition: {dept_attrition['Human Resources']}%
            - Highest Risk Role: {top_roles.index[0]} ({top_roles.iloc[0]}%)

            Write these sections:
            1. Executive Summary (3 sentences)
            2. Key Findings (4 bullet points with numbers)
            3. Root Causes (3 bullet points)
            4. Recommendations (3 bullet points)
            5. Conclusion (2 sentences)

            Use professional HR language.
            """
        }
    ],
    max_tokens=1200
)

ai_report = response.choices[0].message.content
print("✅ AI insights generated")

# ---- Parse AI Report into Sections ----
def extract_section(text, section_name):
    lines   = text.split('\n')
    content = []
    capture = False

    for line in lines:
        if section_name.lower() in line.lower():
            capture = True
            continue
        if capture:
            if any(s in line for s in [
                '2.', '3.', '4.', '5.',
                'Key Finding', 'Root Cause',
                'Recommendation', 'Conclusion',
                'Executive'
            ]) and section_name not in line:
                break
            if line.strip():
                content.append(line.strip())
    return '\n'.join(content)

# ---- Generate Word Document ----
print("📄 Generating Word Report...")

doc = Document()

# Title
title      = doc.add_heading('HR Attrition Analysis Report', 0)
title.runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

# Date and company info
doc.add_paragraph(
    f"Report Date: {datetime.now().strftime('%B %d, %Y')}"
)
doc.add_paragraph(f"Total Employees Analysed: {len(df)}")
doc.add_paragraph(f"Overall Attrition Rate: {attrition_rate}%")
doc.add_paragraph("")

# Key Metrics Table
doc.add_heading('Key Metrics Summary', level=1)
table = doc.add_table(rows=1, cols=2)
table.style = 'Table Grid'

# Header row
header         = table.rows[0].cells
header[0].text = 'Metric'
header[1].text = 'Value'

# Data rows
metrics = [
    ('Total Employees',          str(len(df))),
    ('Employees Left',           str(df['Attrition'].value_counts()['Yes'])),
    ('Employees Stayed',         str(df['Attrition'].value_counts()['No'])),
    ('Overall Attrition Rate',   f'{attrition_rate}%'),
    ('Avg Income - Left',        f'${avg_income_left:,}'),
    ('Avg Income - Stayed',      f'${avg_income_stayed:,}'),
    ('Overtime Attrition Rate',  f"{overtime_attrition['Yes']}%"),
    ('Sales Dept Attrition',     f"{dept_attrition['Sales']}%"),
    ('R&D Dept Attrition',       f"{dept_attrition['Research & Development']}%"),
    ('Highest Risk Job Role',    f"{top_roles.index[0]}")
]

for metric, value in metrics:
    row              = table.add_row().cells
    row[0].text      = metric
    row[1].text      = value

doc.add_paragraph("")

# AI Generated Sections
doc.add_heading('AI-Generated Executive Analysis', level=1)
doc.add_paragraph(ai_report)

# Department Breakdown
doc.add_heading('Department Attrition Breakdown', level=1)
dept_table        = doc.add_table(rows=1, cols=2)
dept_table.style  = 'Table Grid'

dept_header         = dept_table.rows[0].cells
dept_header[0].text = 'Department'
dept_header[1].text = 'Attrition Rate'

for dept, rate in dept_attrition.items():
    row         = dept_table.add_row().cells
    row[0].text = dept
    row[1].text = f'{rate}%'

doc.add_paragraph("")

# Top Job Roles
doc.add_heading('Top 5 High-Risk Job Roles', level=1)
role_table        = doc.add_table(rows=1, cols=2)
role_table.style  = 'Table Grid'

role_header         = role_table.rows[0].cells
role_header[0].text = 'Job Role'
role_header[1].text = 'Attrition Rate'

for role, rate in top_roles.items():
    row         = role_table.add_row().cells
    row[0].text = role
    row[1].text = f'{rate}%'

# Save Word Document
doc.save('HR_Attrition_Report.docx')
print("✅ Word report saved: HR_Attrition_Report.docx")

# ---- Generate PowerPoint ----
print("📊 Generating PowerPoint...")

prs    = Presentation()
width  = prs.slide_width
height = prs.slide_height

# ---- Slide 1: Title Slide ----
slide1  = prs.slides.add_slide(prs.slide_layouts[0])
title1  = slide1.shapes.title
sub1    = slide1.placeholders[1]

title1.text    = "HR Attrition Analysis"
sub1.text      = (
    f"AI-Powered Executive Report\n"
    f"{datetime.now().strftime('%B %Y')}\n"
    f"Total Employees: {len(df)}"
)

title1.text_frame.paragraphs[0].runs[0].font.size  = PptPt(36)
title1.text_frame.paragraphs[0].runs[0].font.bold  = True
title1.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGB(
    0x1F, 0x49, 0x7D
)

# ---- Slide 2: Key Metrics ----
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
body2  = slide2.placeholders[1]

title2.text = "Key Metrics At a Glance"
body2.text  = (
    f"📊 Total Employees: {len(df)}\n"
    f"📉 Overall Attrition Rate: {attrition_rate}%\n"
    f"💰 Avg Income Left: ${avg_income_left:,}\n"
    f"💰 Avg Income Stayed: ${avg_income_stayed:,}\n"
    f"⏰ Overtime Attrition: {overtime_attrition['Yes']}%\n"
    f"✅ No Overtime Attrition: {overtime_attrition['No']}%"
)

# ---- Slide 3: Department Analysis ----
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
body3  = slide3.placeholders[1]

title3.text = "Attrition by Department"
dept_text   = ""
for dept, rate in dept_attrition.items():
    dept_text += f"🏢 {dept}: {rate}%\n"
body3.text  = dept_text

# ---- Slide 4: High Risk Job Roles ----
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
body4  = slide4.placeholders[1]

title4.text  = "Top 5 High-Risk Job Roles"
roles_text   = ""
for i, (role, rate) in enumerate(top_roles.items(), 1):
    roles_text += f"{i}. {role}: {rate}%\n"
body4.text   = roles_text

# ---- Slide 5: Recommendations ----
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
body5  = slide5.placeholders[1]

title5.text = "Key Recommendations"
body5.text  = (
    "1. Conduct compensation review\n"
    "   → Left employees earn $2,046 less per month\n\n"
    "2. Address overtime culture\n"
    f"   → Overtime triples attrition risk\n\n"
    "3. Focus on Sales department\n"
    f"   → Highest attrition at "
    f"{dept_attrition['Sales']}%\n\n"
    "4. Retain Sales Representatives\n"
    f"   → {top_roles.index[0]} at "
    f"{top_roles.iloc[0]}% attrition"
)

# ---- Slide 6: Conclusion ----
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
body6  = slide6.placeholders[1]

title6.text = "Conclusion & Next Steps"
body6.text  = (
    "Immediate Actions Required:\n\n"
    "✅ Review compensation for high-risk roles\n"
    "✅ Implement overtime reduction policies\n"
    "✅ Launch Sales dept retention program\n"
    "✅ Conduct employee satisfaction surveys\n"
    "✅ Set up monthly attrition tracking dashboard\n\n"
    f"Target: Reduce attrition from "
    f"{attrition_rate}% to below 10%"
)

# Save PowerPoint
prs.save('HR_Attrition_Presentation.pptx')
print("✅ PowerPoint saved: HR_Attrition_Presentation.pptx")

# ---- Final Summary ----
print("\n" + "=" * 50)
print("🎉 PHASE 3 COMPLETE!")
print("=" * 50)
print("Files generated:")
print("  📄 HR_Attrition_Report.docx")
print("  📊 HR_Attrition_Presentation.pptx")
print("  📝 hr_analysis_report.txt")
print("\nYour AI Automation Pipeline:")
print("  CSV Data → Python Analysis →")
print("  AI Insights → Word + PowerPoint")
print("=" * 50)